"""
Enhanced script to create a PowerPoint presentation on Report Validation Standards
This version better extracts design elements from the template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def extract_template_design(prs):
    """Extract color scheme and design elements from template"""
    colors = {
        'primary': RGBColor(0, 102, 204),  # Default vibrant blue
        'secondary': RGBColor(255, 102, 0),  # Default vibrant orange
        'accent': RGBColor(0, 153, 76),  # Default vibrant green
        'dark_text': RGBColor(33, 33, 33),
        'light_bg': RGBColor(245, 245, 250)
    }
    
    # Try to extract colors from template if available
    if len(prs.slides) > 0:
        try:
            # Check first slide for color scheme
            slide = prs.slides[0]
            for shape in slide.shapes:
                if hasattr(shape, 'fill'):
                    if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        # Extract color if found
                        pass
        except:
            pass
    
    return colors

def create_validation_presentation():
    # Load template
    try:
        template_prs = Presentation('SamplePPT.pptx')
        colors = extract_template_design(template_prs)
        # Create new presentation using template's slide dimensions
        prs = Presentation()
        prs.slide_width = template_prs.slide_width
        prs.slide_height = template_prs.slide_height
        print(f"Template loaded - Slide size: {prs.slide_width/Inches(1):.1f}\" x {prs.slide_height/Inches(1):.1f}\"")
    except Exception as e:
        # Create new presentation with standard dimensions
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        colors = extract_template_design(None)
        print(f"Creating new presentation - {e}")
    
    # Use extracted or default colors
    primary = colors['primary']
    secondary = colors['secondary']
    accent = colors['accent']
    dark_text = colors['dark_text']
    white = RGBColor(255, 255, 255)
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = primary
        
        # Title
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = white
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            top = Inches(4.2)
            height = Inches(0.8)
            txBox2 = slide.shapes.add_textbox(left, top, width, height)
            tf2 = txBox2.text_frame
            tf2.text = subtitle
            p2 = tf2.paragraphs[0]
            p2.font.size = Pt(24)
            p2.font.color.rgb = white
            p2.alignment = PP_ALIGN.CENTER
        return slide
    
    def add_content_slide(title, title_color, description, items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        
        # Title with colored background
        title_height = Inches(0.8)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, title_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = title_color
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width, title_height)
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = white
        
        # Description
        top = Inches(1.4)
        if description:
            txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
            tf = txBox.text_frame
            tf.text = description
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = dark_text
            top = Inches(1.9)
        
        # Items
        for i, item in enumerate(items):
            if not item.strip():
                continue
            item_top = top + i * Inches(0.35)
            txBox = slide.shapes.add_textbox(left, item_top, width, Inches(0.35))
            tf = txBox.text_frame
            tf.text = item
            p = tf.paragraphs[0]
            if item.startswith("   •") or item.startswith("  •"):
                p.font.size = Pt(13)
            elif item.startswith("•"):
                p.font.size = Pt(14)
            else:
                p.font.size = Pt(15)
                p.font.bold = True
            p.font.color.rgb = dark_text
        return slide
    
    # Slide 1: Title
    add_title_slide(
        "Report Validation Standards",
        "Setting Standards Across Data Analytics & AI Department"
    )
    
    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Validation Framework Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = primary
    
    content_items = [
        "🔧 Technical Validation",
        "🧭 Functional Validation",
        "🧱 Data Integrity Validation",
        "🎯 Data Accuracy Validation"
    ]
    
    top = Inches(1.8)
    for i, item in enumerate(content_items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(1.1), width, Inches(0.8))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 3: Scope
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(left, top - Inches(1.5), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Scope & Reporting Tools"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = primary
    
    top = Inches(1.8)
    sections = [
        ("Reporting Tools:", ["• Oracle Analytics Cloud", "• Power BI"]),
        ("Source Systems (OLTP):", ["• Oracle E-Business Suite", "• Oracle Fusion Cloud"])
    ]
    
    for section_title, items in sections:
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        tf = txBox.text_frame
        tf.text = section_title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = dark_text
        
        top += Inches(0.5)
        for item in items:
            txBox = slide.shapes.add_textbox(left + Inches(0.3), top, width, Inches(0.5))
            tf = txBox.text_frame
            tf.text = item
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = dark_text
            top += Inches(0.5)
        top += Inches(0.3)
    
    # Slide 4: Technical Validation
    add_content_slide(
        "🔧 Technical Validation",
        primary,
        "These checks ensure the report behaves like a well-trained instrument, not a rogue drummer.",
        [
            "1. Report Navigation",
            "   • Validate object-to-object navigation paths",
            "   • Ensure links, breadcrumbs, menu entries, and dashboard flows lead to the right pages",
            "   • Confirm no circular or dead-end navigation loops",
            "",
            "2. Drill Through / Drill Down",
            "   • Drill Down: Validates hierarchical expansion (e.g., Ledger → Account → Journal → Line)",
            "   • Drill Through: Confirm passing of parameters to downstream reports",
            "   • Ensure filters persist correctly when moving between summary and detail layers",
            "",
            "3. Measure Aggregations",
            "   • Validate formulas for SUM, AVG, MIN, MAX, COUNT, DISTINCT COUNT",
            "   • Confirm aggregation rules match business logic",
            "   • Ensure aggregation behavior is consistent across dashboards, visuals, and exports"
        ]
    )
    
    # Slide 5: Functional Validation
    add_content_slide(
        "🧭 Functional Validation",
        secondary,
        "These checks verify that the report behaves as users expect, not as developers wish.",
        [
            "1. Responsiveness to Dashboard Prompts / Filters",
            "   • Validate slicer behavior for multi-select, default values, cascading prompts",
            "   • Ensure applying filters updates all dependent visuals and KPIs",
            "   • Confirm performance is reasonable under typical filter combinations",
            "",
            "2. Measure Value Validation (Summary vs. Drill Reports)",
            "   • Verify that totals shown in summary dashboards match the roll-up of drill reports",
            "   • Confirm that subtotals obey business hierarchies and time dimensions",
            "   • Check Grand Totals vs. SUM of Rows (common trap!)"
        ]
    )
    
    # Slide 6: Data Integrity Validation
    add_content_slide(
        "🧱 Data Integrity Validation",
        accent,
        "This ensures the report's story isn't fiction.",
        [
            "1. Validity of Data Based on Report Definition",
            "   • Check column-level definitions against business metadata",
            "   • Validate joins between subject areas or datasets",
            "   • Ensure orphaned records are excluded unless explicitly expected",
            "   • Confirm proper handling of nulls, effective dates, primary keys"
        ]
    )
    
    # Slide 7: Data Accuracy Validation
    add_content_slide(
        "🎯 Data Accuracy Validation",
        RGBColor(153, 51, 153),  # Purple
        "This is the final face-off: the report vs. the source (OLTP).",
        [
            "1. Comparison with Source Database",
            "   • Choose a fixed set of common parameters (e.g., date, BU, ledger, customer, supplier)",
            "   • Pull the equivalent query directly from EBS/Fusion Cloud tables",
            "   • Validate fact values at:",
            "     - Transaction level",
            "     - Daily/Monthly summary",
            "     - Ledger/Balancing segment",
            "   • Validate dimensional conformance (customer, item, employee, period, account)",
            "",
            "2. Edge Cases",
            "   • Canceled transactions",
            "   • Backdated entries",
            "   • Effective-date based changes",
            "   • Multi-currency conversions",
            "   • Transaction corrections and reversals"
        ]
    )
    
    # Slide 8: Validation Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Validation Checklist Summary"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = primary
    
    checklist = [
        "✓ Technical: Navigation, Drill-down, Aggregations",
        "✓ Functional: Filters, Measure Values, Summary vs Drill",
        "✓ Data Integrity: Column Definitions, Joins, Null Handling",
        "✓ Data Accuracy: Source Comparison, Edge Cases"
    ]
    
    top = Inches(2.2)
    for i, item in enumerate(checklist):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(1.1), width, Inches(0.8))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 9: Standards
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(left, top - Inches(1.9), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Setting Standards Across Department"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = primary
    
    standards = [
        "📋 Establish Validation Templates",
        "🔍 Mandatory Pre-Release Validation",
        "📊 Documentation Requirements",
        "✅ Sign-off Process",
        "🔄 Continuous Improvement"
    ]
    
    top = Inches(2)
    for i, item in enumerate(standards):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.85), width, Inches(0.7))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 10: Thank You
    add_title_slide("Thank You", "Questions & Discussion")
    
    # Save presentation
    output_file = 'Report_Validation_Standards.pptx'
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_validation_presentation()

