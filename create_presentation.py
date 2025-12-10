"""
Script to create a PowerPoint presentation on Report Validation Standards
based on the template from SamplePPT.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_validation_presentation():
    # Try to load template, if it fails, create new with modern design
    try:
        prs = Presentation('SamplePPT.pptx')
        # Use the template's slide layout
        print("Template loaded successfully")
    except:
        # Create new presentation with modern design
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        print("Creating new presentation with modern design")
    
    # Define modern vibrant colors
    primary_color = RGBColor(0, 102, 204)  # Vibrant Blue
    secondary_color = RGBColor(255, 102, 0)  # Vibrant Orange
    accent_color = RGBColor(0, 153, 76)  # Vibrant Green
    dark_text = RGBColor(33, 33, 33)
    light_bg = RGBColor(245, 245, 250)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_color
    
    # Title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Report Validation Standards"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    top = Inches(4)
    height = Inches(0.8)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "Setting Standards Across Data Analytics & AI Department"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Overview/Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Validation Framework Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    content_items = [
        "🔧 Technical Validation",
        "🧭 Functional Validation",
        "🧱 Data Integrity Validation",
        "🎯 Data Accuracy Validation"
    ]
    
    top = Inches(1.5)
    for i, item in enumerate(content_items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(1.2), width, Inches(0.8))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 3: Scope
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(left, top - Inches(1.2), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Scope & Reporting Tools"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    # Reporting Tools
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Reporting Tools:"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = dark_text
    
    top = Inches(2.2)
    tools = ["• Oracle Analytics Cloud", "• Power BI"]
    for i, tool in enumerate(tools):
        txBox = slide.shapes.add_textbox(left + Inches(0.3), top + i * Inches(0.6), width, Inches(0.5))
        tf = txBox.text_frame
        tf.text = tool
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = dark_text
    
    # Source Systems
    top = Inches(3.8)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Source Systems (OLTP):"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = dark_text
    
    top = Inches(4.5)
    systems = ["• Oracle E-Business Suite", "• Oracle Fusion Cloud"]
    for i, system in enumerate(systems):
        txBox = slide.shapes.add_textbox(left + Inches(0.3), top + i * Inches(0.6), width, Inches(0.5))
        tf = txBox.text_frame
        tf.text = system
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = dark_text
    
    # Slide 4: Technical Validation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title with background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top - Inches(1.2), width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary_color
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top - Inches(1.1), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "🔧 Technical Validation"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Description
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Ensures the report behaves like a well-trained instrument"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = dark_text
    
    # Items
    items = [
        "1. Report Navigation",
        "   • Validate object-to-object navigation paths",
        "   • Ensure links, breadcrumbs, menu entries work correctly",
        "   • Confirm no circular or dead-end navigation loops",
        "",
        "2. Drill Through / Drill Down",
        "   • Drill Down: Hierarchical expansion validation",
        "   • Drill Through: Parameter passing to downstream reports",
        "   • Ensure filters persist correctly",
        "",
        "3. Measure Aggregations",
        "   • Validate SUM, AVG, MIN, MAX, COUNT, DISTINCT COUNT",
        "   • Confirm aggregation rules match business logic",
        "   • Ensure consistency across dashboards and exports"
    ]
    
    top = Inches(2.2)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.35), width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(14) if item.startswith("   •") else Pt(16)
        p.font.color.rgb = dark_text
        if item and not item.startswith("   "):
            p.font.bold = True
    
    # Slide 5: Functional Validation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title with background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top - Inches(1.2), width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = secondary_color
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top - Inches(1.1), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "🧭 Functional Validation"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Description
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Verifies that the report behaves as users expect"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = dark_text
    
    # Items
    items = [
        "1. Responsiveness to Dashboard Prompts / Filters",
        "   • Validate slicer behavior (multi-select, default values)",
        "   • Ensure filters update all dependent visuals and KPIs",
        "   • Confirm performance under typical filter combinations",
        "",
        "2. Measure Value Validation (Summary vs. Drill Reports)",
        "   • Verify totals in summary match drill report roll-ups",
        "   • Confirm subtotals obey business hierarchies",
        "   • Check Grand Totals vs. SUM of Rows (common trap!)"
    ]
    
    top = Inches(2.2)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.4), width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(14) if item.startswith("   •") else Pt(16)
        p.font.color.rgb = dark_text
        if item and not item.startswith("   "):
            p.font.bold = True
    
    # Slide 6: Data Integrity Validation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title with background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top - Inches(1.2), width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top - Inches(1.1), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "🧱 Data Integrity Validation"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Description
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Ensures the report's story isn't fiction"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = dark_text
    
    # Items
    items = [
        "1. Validity of Data Based on Report Definition",
        "   • Check column-level definitions against business metadata",
        "   • Validate joins between subject areas or datasets",
        "   • Ensure orphaned records are excluded unless expected",
        "   • Confirm proper handling of nulls, effective dates, primary keys"
    ]
    
    top = Inches(2.2)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.4), width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(14) if item.startswith("   •") else Pt(16)
        p.font.color.rgb = dark_text
        if item and not item.startswith("   "):
            p.font.bold = True
    
    # Slide 7: Data Accuracy Validation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title with background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top - Inches(1.2), width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(153, 51, 153)  # Purple
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top - Inches(1.1), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "🎯 Data Accuracy Validation"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Description
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "The final face-off: the report vs. the source (OLTP)"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = dark_text
    
    # Items
    items = [
        "1. Comparison with Source Database",
        "   • Choose fixed set of common parameters",
        "   • Pull equivalent query from EBS/Fusion Cloud tables",
        "   • Validate fact values at transaction, daily/monthly, ledger levels",
        "   • Validate dimensional conformance",
        "",
        "2. Edge Cases",
        "   • Canceled transactions",
        "   • Backdated entries",
        "   • Effective-date based changes",
        "   • Multi-currency conversions",
        "   • Transaction corrections and reversals"
    ]
    
    top = Inches(2.2)
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.35), width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(13) if item.startswith("   •") else Pt(16)
        p.font.color.rgb = dark_text
        if item and not item.startswith("   "):
            p.font.bold = True
    
    # Slide 8: Validation Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(left, top - Inches(1.2), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Validation Checklist Summary"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    checklist = [
        "✓ Technical: Navigation, Drill-down, Aggregations",
        "✓ Functional: Filters, Measure Values, Summary vs Drill",
        "✓ Data Integrity: Column Definitions, Joins, Null Handling",
        "✓ Data Accuracy: Source Comparison, Edge Cases"
    ]
    
    top = Inches(2)
    for i, item in enumerate(checklist):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(1.2), width, Inches(0.8))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 9: Next Steps / Standards
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(left, top - Inches(1.2), width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Setting Standards Across Department"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    standards = [
        "📋 Establish Validation Templates",
        "🔍 Mandatory Pre-Release Validation",
        "📊 Documentation Requirements",
        "✅ Sign-off Process",
        "🔄 Continuous Improvement"
    ]
    
    top = Inches(2)
    for i, item in enumerate(standards):
        txBox = slide.shapes.add_textbox(left, top + i * Inches(0.9), width, Inches(0.7))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = dark_text
        p.font.bold = True
    
    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_color
    
    # Title
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Report_Validation_Standards.pptx')
    print("Presentation created successfully: Report_Validation_Standards.pptx")

if __name__ == "__main__":
    create_validation_presentation()

